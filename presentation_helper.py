import re
import aiohttp
from pydantic import BaseModel
from bucketManager import BucketManager
from pptx import Presentation
import pptx
from PIL import Image
from io import BytesIO
from urllib.parse import unquote
from pathvalidate import replace_symbol, sanitize_filename
from pptx.util import Inches
from typing import List, Dict, Any


async def download_file_async(url: str) -> BytesIO:
    try:
        async with aiohttp.ClientSession() as session:
            async with session.get(url) as response:
                if response.status != 200:
                    return None
                return BytesIO(await response.read())
    except Exception as e:
        # Можно добавить логирование ошибки, если это нужно
        return None


def resize_image_to_fit(new_image_width, new_image_height, width_in_presentation, height_in_presentation):

    # Calculate aspect ratios
    new_image_aspect_ratio = new_image_width / new_image_height
    target_aspect_ratio = width_in_presentation / height_in_presentation

    # Determine the dimension to scale by
    if new_image_aspect_ratio > target_aspect_ratio:
        # New image is wider than target size, scale based on width
        scale_factor =  width_in_presentation / new_image_width
    else:
        # New image is taller than target size, scale based on height
        scale_factor =  height_in_presentation / new_image_height

    # Apply scale factor to both dimensions to maintain aspect ratio
    target_width = int(new_image_width * scale_factor)
    target_height = int(new_image_height * scale_factor)

    return target_width, target_height

def find_media_info_by_shape(shape, currentrId):
    for rel in shape.part.rels.values():
        if rel.rId == currentrId:
            media_name = rel.target_ref.split('/')[-1]
            return media_name
    return None;


def replace_image_in_presentation(prs, media_uniq_name, new_image_stream):
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.shape_type == 13:  # Check if shape is an image
                media_name = find_media_info_by_shape(shape, shape._element.blip_rId)
                if media_name and media_name == media_uniq_name:
                    try:
                        # Get dimensions of the original image in the presentation
                        image_part = shape.part.related_part(shape._element.blip_rId)
                        image_blob = image_part.blob
                        #width_in_presentation, heigth_in_presentation = get_image_dimensions(image_blob)
                        width_in_presentation = shape.width.inches * 96  # Convert to pixels
                        height_in_presentation = shape.height.inches * 96  # Convert to pixels
                        
                        # Get dimensions of the new image from the stream
                        with Image.open(new_image_stream) as img:
                            new_image_width, new_image_height = img.size
                            
                       
                            new_width, new_height = resize_image_to_fit(new_image_width, new_image_height, width_in_presentation, height_in_presentation)
                            
    
                            
                            img_resized = img.resize((new_width, new_height), Image.ADAPTIVE)

                            
                            
                            #image_resized to BytesIO
                            new_image_stream_resized = BytesIO()
                            img_resized.save(new_image_stream_resized, format='PNG')
                            new_image_stream_resized.seek(0)

                            #add image to current slide
                            #slide.shapes.add_picture(new_image_stream_resized, Inches(0), Inches(0), width=Inches(new_width/96), height=Inches(new_height/96))

                           
                            #  Устанавливаем новый размер формы
                            shape.width = Inches(new_width/96)
                            shape.height = Inches(new_height/96)

                            #save old shape position
                            left = shape.left
                            top = shape.top
                            #delete old shape
                            sp = shape
                            sp.element.getparent().remove(sp.element)

                            #add new shape
                            slide.shapes.add_picture(new_image_stream_resized, left, top, width=Inches(new_width/96), height=Inches(new_height/96))

                             # Replace the image in the presentation
                            # slide_part, rId = shape.part, shape._element.blip_rId
                            # image_part = slide_part.related_part(rId)
                            # image_part.blob = new_image_stream_resized.getvalue()
                            # new_image_stream_resized.close()
                                
                           
                           

                    except Exception as e:
                        print('Ошибка:', e)
                    break
    return prs


def replace_image_background_in_presentation_withoutResize (prs, media_uniq_name, new_image_stream): 
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.shape_type == 13:  # Check if shape is an image
                media_name = find_media_info_by_shape(shape, shape._element.blip_rId)
                if media_name and media_name == media_uniq_name:
                    try:
                        # Replace the image in the presentation
                        slide_part, rId = shape.part, shape._element.blip_rId
                        image_part = slide_part.related_part(rId)
                        image_part.blob = new_image_stream.getvalue()
                    except Exception as e:
                        print('Ошибка:', e)
                    break
    return prs

class PresentationParams:
    def __init__(self, id: str, presentation: str, replacements: List[Dict[str, Any]], resultFileName: str):
        self.id = id
        self.presentation = presentation
        self.replacements = replacements
        self.resultFileName = resultFileName

    @classmethod
    def from_dict(cls, data: Dict[str, Any]):
        return cls(
            id=data['id'],
            presentation=data['presentation'],
            replacements=data['replacements'],
            resultFileName=data['resultFileName']
        )

async def process_single_presentation(presentationInfo: PresentationParams):
    try:
        imageBucket = 'img'
        presentation_url = presentationInfo.presentation
        bucket = BucketManager()
        file_stream = await download_file_async(presentation_url)
        if not file_stream:
            return ""
        presentation = pptx.Presentation(file_stream)
        result_stream = BytesIO()
        for item in presentationInfo.replacements:
            if item.get('media_unique_name') and item.get('assets_file'):
                if bucket.file_exists(imageBucket, item["assets_file"]):
                    image = bucket.get_object_body_all(imageBucket + '/' + item['assets_file'])
                    if not image:
                        return ""
                    byteImgIO = BytesIO(image)
                    byteImgIO.seek(0)
                    with byteImgIO as image_stream:
                        if item["type"] == "background":
                            result_prs = replace_image_background_in_presentation_withoutResize(presentation, item['media_unique_name'], image_stream)
                        else:
                            result_prs = replace_image_in_presentation(presentation, item['media_unique_name'], image_stream)
                        if not result_prs:
                            return ""
                        else:
                            presentation = result_prs
            else:
                return ""

        presentation.save(result_stream)
        result_stream.seek(0)
        cleanedName = replace_symbol(presentationInfo.resultFileName) + '.pptx'
        resultName = re.sub(r'[^\x00-\x7f]', r'', cleanedName)
        folder = 'tests/'
        return bucket.saveFileToFolderAndGetPublicUrl(resultName, folder, result_stream.read())

    except Exception as e:
        print(f"Error: {e.detail}")
        return ""

