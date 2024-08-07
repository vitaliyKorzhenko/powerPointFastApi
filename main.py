import re
import aiohttp
from fastapi import FastAPI, HTTPException
from pydantic import BaseModel
from bucketManager import BucketManager
from pptx import Presentation
import pptx
from PIL import Image
from io import BytesIO
from fastapi.responses import StreamingResponse
from fastapi.middleware.cors import CORSMiddleware
from urllib.parse import unquote
import requests
from pathvalidate import replace_symbol, sanitize_filename
import io
from pptx.util import Inches
from typing import List

#initialize app
app = FastAPI()

#add cors   
origins = [
    "http://pptx.slideedu.com",
    "https://pptx.slideedu.com",
    "http://pptx.slideedu.com:8000",
    "https://pptx.slideedu.com:8000",
    "http://localhost",
    "http://localhost:8000",
    "http://159.203.124.0",
    "http://159.203.124.0:8000",
    
]

app.add_middleware(
    CORSMiddleware,
    allow_origins=origins,
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)



@app.get("/")
def read_root():
    return {"Hello": "API Service is running"}






def print_rel_for_slide(slide):
    relathionships = slide.rels

    for rel in relathionships:
        print("Отношение:")
        print("  Идентификатор отношения:", rel.rId)
        print("  Тип отношения:", rel.reltype)
        print("  Целевой URI:", rel.target_ref)
        print("  Целевой части:", rel.target_part)


def find_media_info_by_shape(shape, currentrId):
    #print('find_media_info_by_shape', currentrId);
    for rel in shape.part.rels.values():
        #print (rel.target_ref)
        #print(rel.rId)
        if rel.rId == currentrId:
            media_name = rel.target_ref.split('/')[-1]
            return media_name
    return None;


def get_slide_background(slide):
    background = slide.background
    #check if background is image
    #get all information about background
    if background:
      print('background', background)
      print('fill type', background.fill.type)
      print('all property fill', dir(background.fill))
      print('FILL IMAGE', dir(background.fill.image))
      
      #print('background type',  background.background_type);
     
     

def find_slide_background(slide):
    for shape in slide.shapes:
        if shape.is_placeholder:
            print('shape is placeholder', shape)
            return shape
    return None



def parse_pptx(prs):
    print('START')
    images = []
    notImages =[];
    
    for slide in prs.slides:
        background_shape = find_slide_background(slide);
        print('background_shape', background_shape)
        for shape in slide.shapes:
            if shape.shape_type == 13:
                currentrId = shape._element.blip_rId;
                mediaName = find_media_info_by_shape(shape, currentrId)
                images.append({
                    "media_name": mediaName,
                    "rId": currentrId,
                    "name": shape.name,
                    "width": shape.width,
                    "height": shape.height,
                    "shape_id": shape.shape_id, 
                    "shape_type": shape.shape_type,
                    "filename": shape.image.filename,
                    "ext": shape.image.ext,
                    "left": shape.left,
                    "top": shape.top,
                    "imageSize": shape.image.size,
                    "sha1": shape.image.sha1,

                })
            else: 
                notImages.append({
                     "name": shape.name,
                    "width": shape.width,
                    "height": shape.height,
                    "shape_id": shape.shape_id, 
                    "shape_type": shape.shape_type,
                })
    return images







from PIL import Image
from io import BytesIO

def get_image_dimensions(image_blob):
    with Image.open(BytesIO(image_blob)) as img:
        return img.size

def is_square(width, height):
    return width == height


def resize_image_to_fit(new_image_width, new_image_height, width_in_presentation, height_in_presentation):
    #print('Ширина и высота нового изображения:', new_image_width, new_image_height)
    #print('Ширина и высота в презентации:', width_in_presentation, height_in_presentation)

    # Calculate aspect ratios
    new_image_aspect_ratio = new_image_width / new_image_height
    target_aspect_ratio = width_in_presentation / height_in_presentation

    # Determine the dimension to scale by
    if new_image_aspect_ratio > target_aspect_ratio:
        # New image is wider than target size, scale based on width
        scale_factor =  width_in_presentation / new_image_width
    else:
        # New image is taller than target size, scale based on height
        scale_factor =  height_in_presentation / new_image_height

    # Apply scale factor to both dimensions to maintain aspect ratio
    target_width = int(new_image_width * scale_factor)
    target_height = int(new_image_height * scale_factor)

    #print('Наши размеры в результате!!!:', target_width, target_height)
    return target_width, target_height



def replace_image_in_presentation(prs, media_uniq_name, new_image_stream):
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.shape_type == 13:  # Check if shape is an image
                media_name = find_media_info_by_shape(shape, shape._element.blip_rId)
                #print('MEDIA NAME', media_name, media_uniq_name);
                if media_name and media_name == media_uniq_name:
                    #print('Найдено изображение для замены:', media_uniq_name)
                    try:
                        # Get dimensions of the original image in the presentation
                        image_part = shape.part.related_part(shape._element.blip_rId)
                        image_blob = image_part.blob
                        #width_in_presentation, heigth_in_presentation = get_image_dimensions(image_blob)
                        width_in_presentation = shape.width.inches * 96  # Convert to pixels
                        height_in_presentation = shape.height.inches * 96  # Convert to pixels
                        #print("ВЫСОТА И ШИРИНА В ПРЕЗЕНТАЦИИ:", width_in_presentation, "x", height_in_presentation)
                        
                        # Get dimensions of the new image from the stream
                        with Image.open(new_image_stream) as img:
                            new_image_width, new_image_height = img.size
                            #print("ВЫСОТА И ШИРИНА НОВОГО ИЗОБРАЖЕНИЯ!!:",  new_image_width, "x", new_image_height)
                            
                       
                            new_width, new_height = resize_image_to_fit(new_image_width, new_image_height, width_in_presentation, height_in_presentation)
                            #print("New size after resize:", new_width, "x", new_height)
                            
    
                            
                            img_resized = img.resize((new_width, new_height), Image.ADAPTIVE)

                            
                            
                            #image_resized to BytesIO
                            new_image_stream_resized = BytesIO()
                            img_resized.save(new_image_stream_resized, format='PNG')
                            new_image_stream_resized.seek(0)

                            #add image to current slide
                            #slide.shapes.add_picture(new_image_stream_resized, Inches(0), Inches(0), width=Inches(new_width/96), height=Inches(new_height/96))

                           
                            #  Устанавливаем новый размер формы
                            shape.width = Inches(new_width/96)
                            shape.height = Inches(new_height/96)

                            #save old shape position
                            left = shape.left
                            top = shape.top
                            #delete old shape
                            sp = shape
                            sp.element.getparent().remove(sp.element)

                            #add new shape
                            slide.shapes.add_picture(new_image_stream_resized, left, top, width=Inches(new_width/96), height=Inches(new_height/96))

                             # Replace the image in the presentation
                            # slide_part, rId = shape.part, shape._element.blip_rId
                            # image_part = slide_part.related_part(rId)
                            # image_part.blob = new_image_stream_resized.getvalue()
                            # new_image_stream_resized.close()
                                
                           
                           

                            #print('Изображение успешно заменено')
                    except Exception as e:
                        print('Ошибка:', e)
                    break
    return prs




def replace_image_background_in_presentation_withoutResize (prs, media_uniq_name, new_image_stream): 
    #print('#REPLACE IMAGE BACKGROUND IN PRESENTATION WITHOUT RESIZE');
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.shape_type == 13:  # Check if shape is an image
                media_name = find_media_info_by_shape(shape, shape._element.blip_rId)
                #print('MEDIA NAME', media_name, media_uniq_name);
                if media_name and media_name == media_uniq_name:
                   # print('Найдено изображение для замены:', media_uniq_name)
                    try:
                        # Replace the image in the presentation
                        slide_part, rId = shape.part, shape._element.blip_rId
                        image_part = slide_part.related_part(rId)
                        image_part.blob = new_image_stream.getvalue()
                        #print('Изображение успешно заменено')
                    except Exception as e:
                        print('Ошибка:', e)
                    break
    return prs






@app.post("/presentation/parsing")
async def parse_presentation(name: str):
    print('START')
    bucket = BucketManager()
    if (bucket.file_exists('pptx', name)):
        file = bucket.getObjectBody('pptx/' + name)
        file_stream = BytesIO(file)
        presentation = pptx.Presentation(file_stream)
        images = parse_pptx(presentation)
        return images
    else:
        return HTTPException(status_code=404, detail="File not found")


#class for body params
class PresentationParseParams(BaseModel):
    presentation: str
    #replacemnt array object media_unique_name: string, assets_file: string

#parisng by link
@app.post("/presentation/parsingByUrl")
async def parse_presentation_by_url(param: PresentationParseParams):
    print('START')
    file_stream = download_file(param.presentation)
    if not file_stream:
        return HTTPException(status_code=404, detail="Presentation file not found")
    presentation = pptx.Presentation(file_stream)
    images = parse_pptx(presentation)
    return images
    
@app.get("/presentation/getTemplates")
async def get_templates():
    bucket = BucketManager()
    templates = bucket.get_files('pptx')
    return templates



#images pptx get
@app.get("/presentation/getImages")
async def get_templates():
    bucket = BucketManager()
    templates = bucket.get_files('img')
    return templates

#get files from results (folder results) 
@app.get("/presentation/getResults")
async def get_results():
    bucket = BucketManager()
    results = bucket.get_files('pptx')
    return results



#get thumbnails from thumbnails folder
@app.get("/presentation/getThumbnails")
async def get_thumbnails():
    bucket = BucketManager()
    thumbnails = bucket.get_files('thumbnails')
    return thumbnails

#class for body params
class PresentationParams(BaseModel):
    presentation: str
    #replacemnt array object media_unique_name: string, assets_file: string
    replacements: list
#replace image in presentation and save it to results (image get from pictures folder)
    resultFileName: str
    


class ThumbnailsParams(BaseModel):
    presentation: str



def download_file(url):
    """
    Downloads a file from the given URL and saves it with the specified filename.
    
    :param url: The URL of the file to download.
    :param filename: The name under which the file will be saved.
    :return: The file stream if successful, or None in case of error.
    """
    try:
        # Request the file content from the URL
        response = requests.get(url)

        if response.status_code == 200:
            # Save the file content
            file_stream = BytesIO(response.content)
            return file_stream
        else:
            print("Error downloading file. Error code:", response.status_code)
            return None

    except Exception as e:
        print("Error:", e)
        return None


                                   
@app.post("/presentation/generateNewPresentationUseUrl")
async def generateNewPresentationUseUrl(presentationInfo: PresentationParams):
    print('BODY', presentationInfo)
    imageBucket = 'img'
    presentation_url = presentationInfo.presentation
    bucket = BucketManager()
    file_stream = download_file(presentation_url)
    if not file_stream:
        return HTTPException(status_code=404, detail="Presentation file not found")
    #print("fileSTREAM NOT NONE");
    presentation = pptx.Presentation(file_stream)
    #images = parse_pptx(presentation);
    #print("IMAGES IN PR", images);
    result_stream = BytesIO()
    for item in presentationInfo.replacements:
            if item.get('media_unique_name') and item.get('assets_file'):
                #print("GO GO GO", item['media_unique_name'], item['assets_file'])
                if bucket.file_exists(imageBucket, item["assets_file"]):
                    image = bucket.getObjectBody(imageBucket + '/' + item['assets_file'])
                    if not image:
                        return HTTPException(status_code=404, detail="Image not found")
                    byteImgIO = BytesIO(image)
                    byteImgIO.seek(0)
                    with byteImgIO as image_stream:
                        if item["type"] == "background":
                            result_prs = replace_image_background_in_presentation_withoutResize(presentation, item['media_unique_name'], image_stream)
                        else:    
                            result_prs = replace_image_in_presentation(presentation, item['media_unique_name'], image_stream)
                        if not result_prs:
                            return HTTPException(status_code=404, detail="Image not found")
                        else:
                            # Update presentation
                            presentation = result_prs
                            #delete byteImgIO
                    
            else:
                return HTTPException(status_code=404, detail="Invalid replacement item")

        # Save updated presentation to results folder
    presentation.save(result_stream)
    result_stream.seek(0)
    cleanedName = replace_symbol(presentationInfo.resultFileName) + '.pptx'
    resultName = re.sub(r'[^\x00-\x7f]',r'', cleanedName) 
    #print("RESULT CLEANED", resultName);

    return StreamingResponse(BytesIO(result_stream.read()), media_type='application/vnd.openxmlformats-officedocument.presentationml.presentation', headers={'Content-Disposition': f'attachment; filename="{resultName}"'})


                                 
@app.post("/presentation/loadNewPresentationToBucketUseUrl")
async def loadNewPresentationToBucketUseUrl(presentationInfo: PresentationParams):
    print('BODY', presentationInfo)
    imageBucket = 'img'
    presentation_url = presentationInfo.presentation
    bucket = BucketManager()
    file_stream = download_file(presentation_url)
    if not file_stream:
        return HTTPException(status_code=404, detail="Presentation file not found")
    #print("fileSTREAM NOT NONE");
    presentation = pptx.Presentation(file_stream)
    #images = parse_pptx(presentation);
    #print("IMAGES IN PR", images);
    result_stream = BytesIO()
    for item in presentationInfo.replacements:
            if item.get('media_unique_name') and item.get('assets_file'):
                #print("GO GO GO", item['media_unique_name'], item['assets_file'])
                if bucket.file_exists(imageBucket, item["assets_file"]):
                    image = bucket.getObjectBody(imageBucket + '/' + item['assets_file'])
                    if not image:
                        return HTTPException(status_code=404, detail="Image not found")
                    byteImgIO = BytesIO(image)
                    byteImgIO.seek(0)
                    with byteImgIO as image_stream:
                        if item["type"] == "background":
                            result_prs = replace_image_background_in_presentation_withoutResize(presentation, item['media_unique_name'], image_stream)
                        else:    
                            result_prs = replace_image_in_presentation(presentation, item['media_unique_name'], image_stream)
                        if not result_prs:
                            return HTTPException(status_code=404, detail="Image not found")
                        else:
                            # Update presentation
                            presentation = result_prs
                            #delete byteImgIO
                    
            else:
                return HTTPException(status_code=404, detail="Invalid replacement item")

        # Save updated presentation to results folder
    presentation.save(result_stream)
    result_stream.seek(0)
    cleanedName = replace_symbol(presentationInfo.resultFileName) + '.pptx'
    resultName = re.sub(r'[^\x00-\x7f]',r'', cleanedName) 
    folder = 'tests/'
    return bucket.saveFileToFolderAndGetPublicUrl(resultName, folder, result_stream.read())




class PresentationsParams(BaseModel):
    presentations: List[PresentationParams]


async def download_file_async(url: str) -> BytesIO:
    print('DOWNLOAD FILE ASYNC', url)
    async with aiohttp.ClientSession() as session:
        async with session.get(url) as response:
            if response.status != 200:
                print('ERROR DOWNLOAD FILE ASYNC', response.status)
                raise HTTPException(status_code=404, detail="Presentation file not found")
            return BytesIO(await response.read())

async def process_single_presentation(presentationInfo: PresentationParams):
    try:
        print('BODY', presentationInfo)
        imageBucket = 'img'
        presentation_url = presentationInfo.presentation
        bucket = BucketManager()
        file_stream = await download_file_async(presentation_url)
        if not file_stream:
            raise HTTPException(status_code=404, detail="Presentation file not found")

        presentation = pptx.Presentation(file_stream)
        result_stream = BytesIO()
        print('FIND PRESENTATION', presentationInfo)
        for item in presentationInfo.replacements:
            if item.get('media_unique_name') and item.get('assets_file'):
                if bucket.file_exists(imageBucket, item["assets_file"]):
                    image = bucket.getObjectBody(imageBucket + '/' + item['assets_file'])
                    if not image:
                        raise HTTPException(status_code=404, detail="Image not found")
                    byteImgIO = BytesIO(image)
                    byteImgIO.seek(0)
                    with byteImgIO as image_stream:
                        if item["type"] == "background":
                            result_prs = replace_image_background_in_presentation_withoutResize(presentation, item['media_unique_name'], image_stream)
                        else:
                            result_prs = replace_image_in_presentation(presentation, item['media_unique_name'], image_stream)
                        if not result_prs:
                            raise HTTPException(status_code=404, detail="Image not found")
                        else:
                            presentation = result_prs
            else:
                raise HTTPException(status_code=404, detail="Invalid replacement item")

        presentation.save(result_stream)
        result_stream.seek(0)
        cleanedName = replace_symbol(presentationInfo.resultFileName) + '.pptx'
        resultName = re.sub(r'[^\x00-\x7f]', r'', cleanedName)
        folder = 'tests/'
        return bucket.saveFileToFolderAndGetPublicUrl(resultName, folder, result_stream.read())

    except HTTPException as e:
        print(f"Error: {e.detail}")
        return ""
    except Exception as e:
        print(f"Unexpected error proccess single!: {e}")
        return ""






@app.post("/presentation/loadListPresentationsToBucketUseUrl")
async def loadListPresentationsToBucketUseUrl(presentationsInfo: PresentationsParams):
    results = []
    for presentationInfo in presentationsInfo.presentations:
        try:
            public_url = await process_single_presentation(presentationInfo)
            results.append({"presentation": presentationInfo.resultFileName, "url": public_url})
        except HTTPException as e:
            results.append({"presentation": presentationInfo.resultFileName, "error": e.detail})
    
    return results




#add model for body params
class ChangePackgroundParams(BaseModel):
    presentation: str
    #replacemnt array object media_unique_name: string, assets_file: string
    background: str
    oldImage: str
    newImage: str

@app.post("presentation/changeBackgroundInSlideAndReturnPresentation")
async def changeBackgroundInSlideAndReturnPresentation (presentationInfo: ChangePackgroundParams):
    print('START REPLACE IMAGE', presentationInfo)
    templatesBucket = 'pptx'
    imageBucket = 'img'
    presentation_name = presentationInfo.presentation
    bucket = BucketManager()
    
    if bucket.file_exists(templatesBucket, presentation_name):
        file = bucket.getObjectBody(templatesBucket + '/' + presentation_name)
        file_stream = BytesIO(file)
        presentation = pptx.Presentation(file_stream)
        newName = presentationInfo.resultFileName + '.pptx'

        result_stream = BytesIO()  
        presentation.save(result_stream)
        result_stream.seek(0)
        cleanedName = replace_symbol(presentationInfo.resultFileName) + '.pptx'
        resultName = re.sub(r'[^\x00-\x7f]',r'', cleanedName) 
        #print("RESULT CLEANED", resultName);
        return StreamingResponse(BytesIO(result_stream.read()), media_type='application/vnd.openxmlformats-officedocument.presentationml.presentation', headers={'Content-Disposition': f'attachment; filename="{newName}"'})



    else:
        return HTTPException(status_code=404, detail="Presentation file not found")


#set background to slides prs




#return presentation by name
@app.get("/presentation/getPresentation")
async def get_presentation(name: str):
    bucket = BucketManager()
    if (bucket.file_exists('pptx', name)):
        file = bucket.getObjectBody('pptx/' + name)
        return StreamingResponse(BytesIO(file), media_type='application/vnd.openxmlformats-officedocument.presentationml.presentation', headers={'Content-Disposition': f'attachment; filename="{name}"'})
    else:
        return HTTPException(status_code=404, detail="File not found")